import os
import re
from typing import Optional
from pathlib import Path
import tempfile
import shutil


class OfficeConverter:
    """
    Office Converter - 办公文档转换服务
    安全地转换各种办公文档格式
    防止 XXE 和其他 XML 相关攻击
    """

    def __init__(self):
        self.temp_dir = Path(tempfile.mkdtemp(prefix="office_convert_"))

    def _sanitize_xml(self, xml_content: str) -> str:
        """
        清理 XML 内容，移除可能导致 XXE 的实体定义和外部引用
        """
        # 移除 DOCTYPE 声明（可能包含实体定义）
        xml_content = re.sub(
            r'<!DOCTYPE\s+[^>]*\[[^\]]*\]>',
            '',
            xml_content,
            flags=re.DOTALL | re.IGNORECASE
        )
        # 移除简单的 DOCTYPE 声明
        xml_content = re.sub(
            r'<!DOCTYPE\s+[^>]*>',
            '',
            xml_content,
            flags=re.IGNORECASE
        )
        # 移除外部实体引用
        xml_content = re.sub(
            r'<!ENTITY\s+[^>]*SYSTEM\s+[^>]*>',
            '',
            xml_content,
            flags=re.IGNORECASE
        )
        xml_content = re.sub(
            r'<!ENTITY\s+[^>]*PUBLIC\s+[^>]*>',
            '',
            xml_content,
            flags=re.IGNORECASE
        )
        return xml_content

    def _create_secure_xml_parser(self):
        """创建安全的 XML 解析器（禁用外部实体）"""
        try:
            from lxml import etree
            parser = etree.XMLParser(
                resolve_entities=False,
                no_network=True,
                load_dtd=False,
            )
            return parser
        except ImportError:
            try:
                from defusedxml import ElementTree as ET
                return ET
            except ImportError:
                import xml.etree.ElementTree as ET
                return ET

    def convert_docx_to_text(self, file_path: str) -> str:
        """将 DOCX 转换为纯文本（安全模式）"""
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs]
            return "\n".join(paragraphs)
        except ImportError:
            # 回退到 zip + XML 解析
            import zipfile
            from xml.etree import ElementTree as ET

            text_parts = []
            with zipfile.ZipFile(file_path, 'r') as zf:
                for item in zf.namelist():
                    if item.endswith('.xml'):
                        xml_content = zf.read(item).decode('utf-8')
                        # 清理 XML 以防止 XXE
                        xml_content = self._sanitize_xml(xml_content)
                        try:
                            root = ET.fromstring(xml_content)
                            for elem in root.iter():
                                if elem.text:
                                    text_parts.append(elem.text)
                        except ET.ParseError:
                            pass
            return "\n".join(text_parts)

    def convert_xlsx_to_text(self, file_path: str) -> str:
        """将 XLSX 转换为纯文本（安全模式）"""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    row_text = [str(cell.value) if cell.value is not None else "" for cell in row]
                    text_parts.append("\t".join(row_text))
            return "\n".join(text_parts)
        except ImportError:
            # 回退到 zip + XML 解析
            import zipfile
            from xml.etree import ElementTree as ET

            text_parts = []
            with zipfile.ZipFile(file_path, 'r') as zf:
                for item in zf.namelist():
                    if item.endswith('.xml'):
                        xml_content = zf.read(item).decode('utf-8')
                        xml_content = self._sanitize_xml(xml_content)
                        try:
                            root = ET.fromstring(xml_content)
                            for elem in root.iter():
                                if elem.text:
                                    text_parts.append(elem.text)
                        except ET.ParseError:
                            pass
            return "\n".join(text_parts)

    def convert_pptx_to_text(self, file_path: str) -> str:
        """将 PPTX 转换为纯文本（安全模式）"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except ImportError:
            # 回退到 zip + XML 解析
            import zipfile
            from xml.etree import ElementTree as ET

            text_parts = []
            with zipfile.ZipFile(file_path, 'r') as zf:
                for item in zf.namelist():
                    if item.endswith('.xml'):
                        xml_content = zf.read(item).decode('utf-8')
                        xml_content = self._sanitize_xml(xml_content)
                        try:
                            root = ET.fromstring(xml_content)
                            for elem in root.iter():
                                if elem.text:
                                    text_parts.append(elem.text)
                        except ET.ParseError:
                            pass
            return "\n".join(text_parts)

    def convert_pdf_to_text(self, file_path: str) -> str:
        """将 PDF 转换为纯文本（安全模式）"""
        try:
            from pdfminer.high_level import extract_text
            return extract_text(file_path)
        except ImportError:
            raise ImportError("pdfminer.six is required for PDF conversion")

    def convert(self, file_path: str, output_format: str = "text") -> str:
        """
        转换文档为指定格式
        """
        path = Path(file_path)
        suffix = path.suffix.lower()

        if suffix == ".docx":
            return self.convert_docx_to_text(file_path)
        elif suffix == ".xlsx":
            return self.convert_xlsx_to_text(file_path)
        elif suffix == ".pptx":
            return self.convert_pptx_to_text(file_path)
        elif suffix == ".pdf":
            return self.convert_pdf_to_text(file_path)
        else:
            raise ValueError(f"Unsupported file format: {suffix}")

    def cleanup(self):
        """清理临时文件"""
        if self.temp_dir.exists():
            shutil.rmtree(self.temp_dir, ignore_errors=True)
